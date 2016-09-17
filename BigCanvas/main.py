import numpy as np
import cv2
import sys


sys.path.append('/home/salar/CURRENT_PROJECTS/FYP2')

from filters import InkFilter
from CamScribbleCore.Filters.PerspectiveCorrection import PerspectiveCorrection

from pptx import Presentation

CAMERA = 1


def addSlide(image,name='output.pptx'):
	"""
	just pass this function an image and it will add the image to the slideshow
	named name
	:param image: the image to save
	:param name: the name of the slideshow to add to
	:return: nothing
	"""

	cv2.imwrite('picture.png',image)

	try:
		prs = Presentation(name)
	except(Exception):
		prs = Presentation()

	blank_layout = prs.slide_layouts[6]


	slide = prs.slides.add_slide(blank_layout)

	pic = slide.shapes.add_picture('picture.png',0, 0)
	prs.slide_height = pic.height
	prs.slide_width = pic.height

	prs.save('output.pptx')

# The objective here is to develop classes which can be used
# from outside openCV, so please do not create any cv2 windows inside
# just apply filters and return the results as if you have no idea
# about the windowing library that is being used.

# The objective here is to have a large canvas
# which will be layed under a smaller one.
# The output of the Chafochen module will be
# rendered in the smaller canvas which always resides
# inside the bounds of the larger canvas.
# The smaller canvas can be moved at any position over the
# the larger canvas via keyboard input.
# The user initially moves the smaller window to a
# position in the larger canvas and starts
# articulating what he desires. When the user wants
# to transfer the contents of the smaller canvas to
# the larger canvas underneath, he can press 'L'.
# Upon doing so the contents of the smaller window get
# transferred to the larger canvas, and the user can
# move the smaller canvas to a new place.
# ALGORITHM:
#	1) construct a canvas having certain height and width
#	2) set the dimensions of the smaller canvas
#	3) read a frame from the Chafochen module and overlay it
#		on the larger canvas
#	4) check if any keys have been press, and take relevant actions
#	5) goto 3.


class Canvas:
	"""
		The purpose of this class is to implement the smaller filter.
		It will take input from our video source, pass it to InkFilter
		to get the binary mask, and draw that mask at a certain position
		in a big picture.

		The big picture will be passed as an argument to the function
		Canvas.draw().
	"""
	def __init__(self,(width,height),source,inkFilter,perspectiveCorrection):
		"""
		:param (width,height):  Initial size of the canvas.
		:param source:          cv2.VideoCapture object
		:param inkFilter:       pre-configured filters.InkFilter object
		:param perspectiveCorrection: pre-configured PerspectiveCorrection object
		"""

		# set the video source
		self.__cam = source

		# set the InkFilter
		self.__filter = inkFilter

		# set the perspective correction object
		self.__perspectiveCorrection = perspectiveCorrection



		self.__width = width
		self.__height = height

		# white colored ink by default
		self.__inkColor = [255,0,255]

		# offset of the canvas in the big picture
		self.__x_offset = 0
		self.__y_offset = 0



		self.__currentFrameShape = (0,0)

	def draw(self,img,borders=True):
		"""

		:param img: The big picture, this should be an RGB image
		:param borders: Weather to draw the borders around this canvas or not
		:return: nothing
		"""

		# read a frame from the source.
		ret,frame = self.__cam.read()



		if(ret == False):
			return False

		cv2.imshow('original', cv2.resize(frame, (300, 300), interpolation=cv2.INTER_AREA))

		# convert frame to gray scale
		frame = cv2.cvtColor(frame,cv2.COLOR_BGR2GRAY)


		# apply perspective correction
		frame = self.__perspectiveCorrection.applyPerspectiveCorrection(frame)



		# apply InkFilter to the image.
		mask = self.__filter.applyFilter(frame)

		# resize the image
		mask = cv2.resize(mask, (self.__width,self.__height), interpolation=cv2.INTER_AREA)



		# draw borders if required
		if (borders == True):
			# set borderColor to the selected color
			# if borders are to be drawn
			borderColor = 255
		else:
			# else set the border color to black
			borderColor = 0

		mask = cv2.copyMakeBorder(mask, 5, 5, 5, 5, cv2.BORDER_CONSTANT, value=borderColor)

		# cv2.imshow('canvas', mask)


		# apply the color to the current frame
		r = np.empty(mask.shape,dtype=np.uint8)
		g = np.empty(mask.shape, dtype=np.uint8)
		b = np.empty(mask.shape, dtype=np.uint8)

		# cv2.imshow("mask",mask)

		r.fill(self.__inkColor[0])
		cv2.bitwise_and(mask,r,r)

		g.fill(self.__inkColor[1])
		cv2.bitwise_and(mask, g, g)

		b.fill(self.__inkColor[2])
		cv2.bitwise_and(mask, b, b)

		color = cv2.merge((b, g, r))

		# cv2.imshow("color mask", color)

		# r = cv2.bitwise_and(r,currentFrame)
		# g = cv2.bitwise_and(g, currentFrame)
		# b = cv2.bitwise_and(b, currentFrame)







		# now check if the canvas can be drawn on the image.
		# exit the application if the canvas is too large for the
		# image.
		if ((img.shape[0] < mask.shape[0]) or (img.shape[1] < mask.shape[1])):
			print "Cannot draw canvas on image, image too small:\n" +"size of canvas: " + str(mask.shape) + "\n"+"size of image: " + str(img.shape);

			# exit safely
			sys.exit(0)

		# store the handle for the current frame in the object
		# it is used when setting the x_offset and the y_offset,
		# in the function self.setPosition()
		self.__currentFrameShape = mask.shape

		# draw canvas on top of the big image.
		s_img = mask

		x_offset = self.__x_offset
		y_offset = self.__y_offset
		# iterate over the first three channels in
		# the image. 0 1 2 : B G R
		# and apply the smaller canvas

		# This mask is the bitwise not of 'mask'
		notMask = cv2.bitwise_not(mask)
		notMask[notMask != 255] = 0

		# cv2.imshow("not mask",notMask)

		for c in range(0,3):

			# ROI is the copy of the region of the big picture that lies underneath
			# this canvas
			ROI = img[y_offset:y_offset+s_img.shape[0], x_offset:x_offset+s_img.shape[1], c].copy()

			colorMask = color[:,:,c]

			# First and the ROI with the not mask
			cv2.bitwise_and(ROI, notMask, ROI)

			# Then or the ROI with the color mask
			cv2.bitwise_or(ROI,colorMask,ROI)

			# assign ROI back to the image.
			img[y_offset:y_offset + s_img.shape[0], x_offset:x_offset + s_img.shape[1], c] = ROI

		return True

	def setWidth(self,width,shape):
		# -ve values are not accepted
		if(width <=0):
			return False


		# calculate the bound of the new size
		newXBoundary = self.__x_offset+width+10
		# If the horizontal boundary is out of bounds
		if( newXBoundary > shape[1]):
			# calculate the new x offset
			newXOffset = self.__x_offset - (newXBoundary - shape[1])

			# if the new offset is -ve, it means that the specified width
			# is out of bounds of the shape.
			if(newXOffset < 0):
				return False

			self.__x_offset = newXOffset



		self.__width = width
		return True

	def setHeight(self, height,shape):
		# -ve values are not accepted.
		if (height <= 0):
			return False

		# calculate the bound of the new size
		newYBoundary = self.__y_offset + height + 10

		# If the vertical boundary is out of bounds
		if (newYBoundary > shape[0]):
			# calculate the newYoffset.
			newYOffset = self.__y_offset - (newYBoundary - shape[0])

			# a -ve newYOffset will tell us that the specified size
			# is out of bounds of 'shape'
			if(newYOffset < 0):
				return False

			# else assign the newYOffset.
			self.__y_offset = newYOffset


		self.__height = height
		return True

	def setColor(self,color):
		self.__inkColor = color

	def setPosition(self,shape,x,y):

		xPosAfter = x + self.__currentFrameShape[1]
		yPosAfter = y + self.__currentFrameShape[0]


		if(x < 0 ):
			x = 0

		if(y < 0):
			y = 0

		if(xPosAfter > shape[1]):
			x = shape[1] - self.__currentFrameShape[1]

		if(yPosAfter > shape[0]):
			y = shape[0] - self.__currentFrameShape[0]



		self.__x_offset = x
		self.__y_offset = y

	def getX(self):
		return self.__x_offset

	def getY(self):
		return self.__y_offset


class BigPicture:

	def __init__(self,(width,height),cam,perspectiveTransform):

		# # pre-configured cv2.VideoCapture object
		# self.__cam = cam
		#
		# # pre-configured PerspectiveCorrection object
		# self.__perspCorrection = perspectiveTransform

		# initialize the bigPicture
		if width<400 or height<400:
			print "Error@BigPicture.__init__(): the width and height cannot be less than 400 pixels"

		self.__bigPicture = np.ndarray((width,height,3),dtype=np.uint8)

		# set the filter object
		self.filter = InkFilter()


		# initialize the canvas
		self.__canvas = Canvas((400,400),cam,self.filter,perspectiveTransform)


		self.__undoStack = []

		self.__redoStack = []

		self.__undoStack.append(self.__bigPicture.copy())







	def getNextFrame(self):
		img = self.__bigPicture.copy()


		# draw canvas with borders
		if (self.__canvas.draw(img, True) == False):
			return False


		return img

	def setPosition(self,x,y):
		self.__canvas.setPosition(self.__bigPicture.shape, x, y)

	def setHeight(self,height):
		self.__canvas.setHeight(height,self.__bigPicture.shape)

	def setWidth(self,width):
		self.__canvas.setWidth(width, self.__bigPicture.shape)

	def setColor(self,color):
		self.__canvas.setColor(color)


	def getX(self):
		return self.__canvas.getX()

	def getY(self):
		return self.__canvas.getY()


	def lock(self):
		if (self.__canvas.draw(self.__bigPicture, False) == False):
			return False

		# push this saved bigPicture to the undo stack
		self.__undoStack.append(self.__bigPicture.copy())

		# empty the redo stack
		if (len(self.__redoStack) > 0):
			self.__redoStack = []

	def undo(self):
		# undo
		if (len(self.__undoStack) > 1):
			self.__redoStack.append(self.__undoStack.pop())
			self.__bigPicture = self.__undoStack[len(self.__undoStack) - 1].copy()

	def redo(self):
		# redo
		if (len(self.__redoStack) > 0):
			self.__undoStack.append(self.__redoStack.pop())
			self.__bigPicture = self.__undoStack[len(self.__undoStack) - 1].copy()

	def createPicture(self,name='picture.png'):
		cv2.imwrite(name,self.__bigPicture)

	def addSlide(self,name='slides.pptx'):

		cv2.imwrite('temp_slide.png', self.__bigPicture)

		try:
			prs = Presentation(name)
		except(Exception):
			prs = Presentation()

		blank_layout = prs.slide_layouts[6]

		slide = prs.slides.add_slide(blank_layout)

		pic = slide.shapes.add_picture('temp_slide.png', 0, 0)
		prs.slide_height = pic.height
		prs.slide_width = pic.height

		prs.save(name)

	def clear(self):

		self.__bigPicture = np.zeros(self.__bigPicture.shape,dtype=np.uint8)



pointIndex = 0




# # define the big picture
# bigPicture = np.zeros((700,700,3),dtype=np.uint8)
#
#
# undo_stack = []
# redo_stack = []
#
# undo_stack.append(bigPicture.copy())

# create the video source object and
# the InkFilter object.


cam = cv2.VideoCapture(CAMERA)
perspectiveCorrection = PerspectiveCorrection()

application = BigPicture((700,700),cam,perspectiveCorrection)


def setup(cam,inkFilter,perspectiveCorrection):

	img = None

	pts = [(0,0),(0,0),(0,0),(0,0)]

	# local function for handling mouse events
	def draw_circle(event, x, y, flags, param):
		global pointIndex

		if event == cv2.EVENT_LBUTTONDBLCLK:
			cv2.circle(img, (x, y), 0, (255, 0, 255), -1)
			pts[pointIndex] = (x, y)
			pointIndex = pointIndex + 1

	# local function for letting the user select points
	def selectFourPoints():
		_, img = cam.read()
		_, img = cam.read()

		print "Please select 4 points, by double clicking on each of them in the order: \n"+"top left, top right, bottom left, bottom right."

		while (pointIndex != 4):

			cv2.imshow('select four corners', img)
			key = cv2.waitKey(20) & 0xFF
			if key == 27:
				return False

		return True

	# create a named window.
	cv2.namedWindow('select four corners')

	# register mouse click callback.
	cv2.setMouseCallback('select four corners', draw_circle)


	# make the user select 4 points
	# and save the 4 points in the array 'pts'
	# exit if the user presses ESC.
	if (selectFourPoints() is False):
		sys.exit(0)


	cv2.destroyWindow('select four corners')

	# configure the int filter
	perspectiveCorrection.setSurfaceCorners(pts[0],pts[1],pts[2],pts[3])

	# perspectiveCorrection.setOutputHeight(500)
	# perspectiveCorrection.setOutputWidth(500)

	def threshChanged(val):
		application.filter.setAdaptiveFilterThreshold(val)


	def aThreshChanged(val):
		application.filter.setAdaptiveFilterSize(val)

	# def medianChanged(val):
	# 	application.filter.setMedianFilterSize(val)

	# def gaussianChanged(val):
	# 	application.filter.setGaussianFilterSize(val)

	# def clipLimit(val):
	# 	application.filter.setCLAHEClipLimit(val)

	# def tileGridSize(val):
	# 	application.filter.setCLAHEGridSize(val)

	# def bgThreshold(val):
	# 	application.filter.setBGThreshold(val)

	cv2.namedWindow('internal controls')
	cv2.createTrackbar('threshold', 'internal controls', 0, 255, threshChanged)
	cv2.createTrackbar('adaptive filter', 'internal controls', 3, 25, aThreshChanged)

	# cv2.createTrackbar('median', 'internal controls', 0, 30, medianChanged)
	# cv2.createTrackbar('gaussian', 'internal controls', 0, 30, gaussianChanged)

	# cv2.createTrackbar('clip limit', 'internal controls', 0, 255, clipLimit)
	# cv2.createTrackbar('grid size', 'internal controls', 0, 50, tileGridSize)

	# cv2.createTrackbar('BG thresh', 'internal controls', 0, 255, bgThreshold)


# confiure the inkFilter object
setup(cam,application.filter,perspectiveCorrection)




COLOR = [255,255,255]

def redChanged(val):
	COLOR[0] = val
	application.setColor(COLOR)

def greenChanged(val):
	COLOR[1] = val
	application.setColor(COLOR)


def blueChanged(val):
	COLOR[2] = val
	application.setColor(COLOR)

def hueChanged(val):
	global COLOR

	HSV = np.array([[[val,255,255]]],dtype=np.uint8)

	RGB = cv2.cvtColor(HSV,cv2.COLOR_HSV2BGR)

	COLOR = RGB[0][0]

	application.setColor(COLOR)

# cv2.createTrackbar('R','color',0,255,redChanged)
# cv2.createTrackbar('G','color',0,255,greenChanged)
# cv2.createTrackbar('B','color',0,255,blueChanged)

cv2.createTrackbar('Hue','internal controls',0,180,hueChanged)


cv2.namedWindow('big picture')


########################MOUSE CALLBACKS########################

DRAG_ENABLED = False
def mouse_controller(event,x,y,flags,param):
	global DRAG_ENABLED
	global canv
	global bigPicture
	if event == cv2.EVENT_LBUTTONDOWN:
		DRAG_ENABLED = True

	# check to see if the left mouse button was released
	elif event == cv2.EVENT_LBUTTONUP:
		DRAG_ENABLED = False

	elif event == cv2.EVENT_MOUSEMOVE:
		if(DRAG_ENABLED):
			application.setPosition(x,y)

cv2.setMouseCallback('big picture',mouse_controller)
#################################################
#################### SIZING #####################
def heightChanged(val):
	application.setHeight(val)
def widthChanged(val):
	application.setWidth(val)

cv2.createTrackbar('height', 'internal controls', 0, 700-10, heightChanged)
cv2.createTrackbar('width', 'internal controls', 0, 700-10, widthChanged)

#################################################

PAUSED = False

SLIDE = 0
while(1):
	# if (PAUSED == False):
	# 	img = bigPicture.copy()
	#
	#
	# if(canv.draw(img,True) == False):
	# 	break
	#
	# cv2.imshow("big picture",img)

	img = application.getNextFrame()

	cv2.imshow("big picture",img)


	key = cv2.waitKey(1) & 0xFF

	if key == 27:
		break

	# left
	elif key == ord('a'):
		application.setPosition(application.getX()-5,application.getY())

	# up
	elif key == ord('w'):
		canv.setPosition(application.getX(), application.getY() -5)

	# right
	elif key == ord('d'):
		canv.setPosition(application.getX() + 5, application.getY())

	# down
	elif key == ord('s'):
		canv.setPosition(application.getX(), application.getY() + 5)

	elif key == ord('l'):
		# if(canv.draw(bigPicture,False) == False):
		# 	break
		#
		# # push this saved bigPicture to the undo stack
		# undo_stack.append(bigPicture.copy())
		#
		# # empty the redo stack
		# if(len(redo_stack) > 0):
		# 	redo_stack = []
		# # bigPicture = img
		application.lock()

	elif key == ord('c'):
		# addSlide(bigPicture)
		# # cv2.imwrite(str(SLIDE)+".png",bigPicture)
		# # SLIDE += 1
		application.addSlide('Slides.pptx')

	elif key == ord('m'):
		application.clear()

	elif key == ord('p'):
		PAUSED = not PAUSED

	elif key == ord('u'):
		application.undo()

	elif key == ord('r'):
		application.redo()

	elif key == ord('g'):
		application.filter.adaptive_gaussian = not application.filter.adaptive_gaussian


cv2.destroyAllWindows()